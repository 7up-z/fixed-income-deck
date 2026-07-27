from __future__ import annotations

import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT.parents[1] / "upload" / "固收业务展业研究汇报材料20260713.pptx"
OUTPUT = ROOT / "public" / "data" / "slides.json"
CANVAS_W = 1672
CANVAS_H = 941


def align_name(value):
    return {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
    }.get(value, "left")


def extract_text_boxes(slide, sx, sy):
    boxes = []
    for index, shape in enumerate(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip()
        if not text:
            continue
        first_paragraph = shape.text_frame.paragraphs[0]
        runs = first_paragraph.runs
        font_pt = runs[0].font.size.pt if runs and runs[0].font.size else 16
        x = max(0, round(shape.left * sx, 2))
        y = max(0, round(shape.top * sy, 2))
        width = min(round(shape.width * sx, 2), CANVAS_W - x)
        height = min(round(shape.height * sy, 2), CANVAS_H - y)
        boxes.append(
            {
                "key": f"text-{index:03d}",
                "label": text.replace("\n", " ")[:32],
                "text": text,
                "x": x,
                "y": y,
                "w": width,
                "h": height,
                "fontSize": round(font_pt * 1.333, 1),
                "align": align_name(first_paragraph.alignment),
            }
        )
    return boxes


def main():
    deck = Presentation(PPTX)
    sx = CANVAS_W / deck.slide_width
    sy = CANVAS_H / deck.slide_height
    slides = []
    for number, slide in enumerate(deck.slides, 1):
        slides.append(
            {
                "id": f"{number:02d}",
                "textBoxes": extract_text_boxes(slide, sx, sy),
            }
        )
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    OUTPUT.write_text(
        json.dumps({"slides": slides}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(f"wrote {OUTPUT} with {len(slides)} slides")


if __name__ == "__main__":
    main()
